from flask import Flask, request, jsonify
from flask_cors import CORS
from werkzeug.utils import secure_filename
from pypdf import PdfReader
from pptx import Presentation
from docx import Document
from PIL import Image
import pytesseract

from qwen_client import generate_study_material

app = Flask(__name__)
CORS(app)

# Reject uploads bigger than 10MB
app.config["MAX_CONTENT_LENGTH"] = 10 * 1024 * 1024

ALLOWED_EXTENSIONS = {"pdf", "pptx", "docx", "png", "jpg", "jpeg"}
IMAGE_EXTENSIONS = {"png", "jpg", "jpeg"}
MAX_EXTRACTED_CHARS = 12000


def allowed_file(filename):
    return "." in filename and filename.rsplit(".", 1)[1].lower() in ALLOWED_EXTENSIONS


def extract_pdf_text(file_stream):
    reader = PdfReader(file_stream)
    pages = []
    for page in reader.pages:
        text = page.extract_text()
        if text:
            pages.append(text)
    return "\n".join(pages)


def extract_pptx_text(file_stream):
    presentation = Presentation(file_stream)
    lines = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                line = "".join(run.text for run in paragraph.runs)
                if line.strip():
                    lines.append(line)
    return "\n".join(lines)


def extract_docx_text(file_stream):
    document = Document(file_stream)
    paragraphs = [p.text for p in document.paragraphs if p.text.strip()]

    # Also pull text out of any tables in the document
    for table in document.tables:
        for row in table.rows:
            for cell in row.cells:
                if cell.text.strip():
                    paragraphs.append(cell.text)

    return "\n".join(paragraphs)


def extract_image_text(file_stream):
    image = Image.open(file_stream)
    return pytesseract.image_to_string(image)


@app.route("/")
def home():
    return {"message": "StudyFlow AI backend is running"}


@app.route("/extract", methods=["POST"])
def extract():
    if "file" not in request.files:
        return jsonify({"error": "No file uploaded."}), 400

    file = request.files["file"]

    if file.filename == "":
        return jsonify({"error": "No file selected."}), 400

    if not allowed_file(file.filename):
        return jsonify({"error": "Only PDF and PPTX files are supported."}), 400

    filename = secure_filename(file.filename)
    extension = filename.rsplit(".", 1)[1].lower()

    try:
        if extension == "pdf":
            text = extract_pdf_text(file)
        elif extension == "pptx":
            text = extract_pptx_text(file)
        elif extension == "docx":
            text = extract_docx_text(file)
        elif extension in IMAGE_EXTENSIONS:
            text = extract_image_text(file)
        else:
            return jsonify({"error": "Unsupported file type."}), 400
    except Exception:
        return jsonify({
            "error": "Couldn't read that file. Make sure it isn't corrupted or password-protected."
        }), 400

    text = text.strip()

    if not text:
        return jsonify({"error": "No readable text was found in that file."}), 400

    if len(text) > MAX_EXTRACTED_CHARS:
        text = text[:MAX_EXTRACTED_CHARS]

    return jsonify({"text": text})


@app.errorhandler(413)
def file_too_large(e):
    return jsonify({"error": "File is too large. Please upload something under 10MB."}), 413


@app.route("/generate", methods=["POST"])
def generate():
    data = request.get_json()
    course_notes = data.get("course_notes", "")

    if not course_notes.strip():
        return jsonify({"error": "Please provide course notes."}), 400

    result = generate_study_material(course_notes)

    return jsonify({"result": result})


if __name__ == "__main__":
    app.run(debug=True, port=5001)